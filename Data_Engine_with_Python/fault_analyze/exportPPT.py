from pptx import Presentation
import  os
from pptx import Presentation
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches
###################这个从程序用于将报警导出到PPT，需要一个PPT模板的文件。。 姜帆  2020年7月######
filename = input('输入保存地址!')
save='E:\DATA_ENGIN'+'\\'+filename+'\\'
##################*******************######################
prs = Presentation('E:\DATA_ENGIN'+'\\'+'bk.pptx')
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = filename+"故障数据分析"
subtitle.text = "PFA3P-120JPH"
####2
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
subtitle.text = "Process设备数据"
textbox= slide.shapes.add_textbox(left=Inches(3),
                                    top=Inches(3),
                                    width=Inches(8),
                                    height=Inches(2)
                                   )
tf = textbox.text_frame
para = tf.add_paragraph()    # 添加段落
para.text = '故障词云'
para.alignment = PP_ALIGN.CENTER
para = tf.add_paragraph() 
font = para.font
font.size = Pt(36)    # 大小
font.name = '华文彩云'    # 字体
font.bold = True    # 加粗
font.italic = True  # 倾斜
font.color.rgb = RGBColor(225, 225, 0)  # 黄色
imagname=save+'process故障.png'
""" pic = slide.shapes.add_picture(image_file=imagname,
                           left=Inches(0.5),
                           top=Inches(2.5),
                           width=Inches(3),
                           height=Inches(3)
                          ) """
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = filename+"故障数据分析"
img_path9 = save+'process每小时曲线bar.jpg' 
img_path4= save+'process故障报警分型次数bar.jpg' 
img_path1 = save+'process每日故障时间bar.jpg' 
img_path0 =save+'process每日故障次数bar.jpg'
img_path3 = save+'process设备条线故障次数bar.jpg' 
img_path2 = save+'process故障时长bar.jpg' 
img_path5 = save+'process设备元器件故障次数bar.jpg' 
img_path6 = save+'process设备种类分型次数bar.jpg' 
img_path7 = save+'process设备TOP10频次bar.jpg' 
img_path8 = save+'process设备TOP10时长bar.jpg' 

 #图片名称一定要对

top = Inches(1.5)

left = Inches(0.5)
height = Inches(3)

blank_slide_layout = prs.slide_layouts[2]
title = slide.shapes.title
subtitle = slide.placeholders[1]
title_slide_layout = prs.slide_layouts[0]
title.text = filename+"故障数据分析"
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path9, left, top)
slide = prs.slides.add_slide(blank_slide_layout)  
pic = slide.shapes.add_picture(img_path0, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path1, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path7, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path8, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path2, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path3, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path4, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path5, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path6, left, top)
# slide = prs.slides.add_slide(blank_slide_layout)
# pic = slide.shapes.add_picture(img_path7, left, top)

##################################################

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title


subtitle = slide.placeholders[1]
title.text = filename+"故障数据分析"
subtitle.text = "Conveyor设备数据"
imagname=save+'conveyor故障.png'
""" pic = slide.shapes.add_picture(image_file=imagname,
                           left=Inches(0.5),
                           top=Inches(2.5),
                           width=Inches(3),
                           height=Inches(3)
                          ) """
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = filename+"故障数据分析"
eqmt = save+'conveyor'
img_path4= eqmt+'故障报警分型次数bar.jpg' 
img_path1 = eqmt+'每日故障时间bar.jpg' 
img_path0 = eqmt+'每日故障次数bar.jpg' 
img_path3 = eqmt+'设备条线故障次数bar.jpg' 
img_path2 = eqmt+'故障时长bar.jpg' 
img_path5 = eqmt+'设备元器件故障次数bar.jpg' 
img_path6 = eqmt+'设备种类分型次数bar.jpg' 
img_path7 = eqmt+'设备TOP10频次bar.jpg' 
img_path8 = eqmt+'设备TOP10时长bar.jpg' 
img_path9 = eqmt+'每小时曲线bar.jpg'
 #图片名称一定要对

top = Inches(1.5)

left = Inches(0.5)
height = Inches(3)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path9, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path0, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path1, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path7, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path8, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path2, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path3, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path4, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path5, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path6, left, top)
# slide = prs.slides.add_slide(blank_slide_layout)
# pic = slide.shapes.add_picture(img_path9, left, top)
#######################################################################


title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title


subtitle = slide.placeholders[1]
title.text = filename+"故障数据分析"
subtitle.text = "机器人设备数据"
imagname=save+'APT故障.png'
""" pic = slide.shapes.add_picture(image_file=imagname,
                           left=Inches(0.5),
                           top=Inches(2.5),
                           width=Inches(3),
                           height=Inches(3)
                          ) """
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = filename+"故障数据分析"
eqmt = save+'APT'
img_path4= eqmt+'故障报警分型次数bar.jpg' 
img_path1 = eqmt+'每日故障时间bar.jpg' 
img_path0 = eqmt+'每日故障次数bar.jpg' 
img_path3 = eqmt+'设备条线故障次数bar.jpg' 
img_path2 = eqmt+'故障时长bar.jpg' 
img_path5 = eqmt+'设备元器件故障次数bar.jpg' 
img_path6 = eqmt+'设备种类分型次数bar.jpg' 
img_path7 = eqmt+'设备TOP10频次bar.jpg' 
img_path8 = eqmt+'设备TOP10时长bar.jpg' 
img_path9 = eqmt+'每小时曲线bar.jpg' 
 #图片名称一定要对

top = Inches(1.5)

left = Inches(0.5)
height = Inches(3)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path9, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path0, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path1, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path7, left, top)
slide = prs.slides.add_slide(blank_slide_layout)   
pic = slide.shapes.add_picture(img_path8, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path2, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path3, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path4, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path5, left, top)
slide = prs.slides.add_slide(blank_slide_layout)
pic = slide.shapes.add_picture(img_path6, left, top)
print('完成！！')
prs.save(save+'故障数据分析报告.pptx')




